"""
Complete Document Translation System
Supports PDF, Word, Excel, and PowerPoint with format preservation
Optimized for macOS but works cross-platform
"""

import os
import sys
import platform
from pathlib import Path
from typing import Optional, Dict, List, Union, Any
from functools import lru_cache
from enum import Enum
import json
import argparse
from datetime import datetime

# Core libraries
try:
    import pymupdf  # PDF support
except ImportError:
    print("Warning: PyMuPDF not installed. PDF support disabled.")
    pymupdf = None

try:
    from pptx import Presentation  # PowerPoint support
except ImportError:
    print("Warning: python-pptx not installed. PowerPoint support disabled.")
    Presentation = None

try:
    # Excel support (openpyxl >= 3.1 no longer ships formula tokenizer)
    from openpyxl import load_workbook, Workbook  # Excel support
    try:
        from openpyxl.formula import Tokenizer, Token  # optional in older versions
    except Exception:
        Tokenizer = None
        Token = None
    try:
        from openpyxl.utils.exceptions import InvalidFileException
    except Exception:
        InvalidFileException = Exception
except ImportError:
    print("Warning: openpyxl not installed. Excel support disabled.")
    load_workbook = None

try:
    from docx import Document  # Word support
except ImportError:
    print("Warning: python-docx not installed. Word support disabled.")
    Document = None

# Translation libraries
try:
    from deep_translator import GoogleTranslator, DeeplTranslator, MicrosoftTranslator
except ImportError:
    print("Error: deep-translator not installed. Please install with: pip install deep-translator")
    sys.exit(1)

# Optional: DeepL official library
try:
    import deepl
except ImportError:
    deepl = None

# For retry logic
try:
    from tenacity import retry, stop_after_attempt, wait_exponential
except ImportError:
    print("Warning: tenacity not installed. Retry logic disabled.")
    retry = lambda *args, **kwargs: lambda func: func


class DocumentType(Enum):
    """Supported document types"""
    PDF = "pdf"
    WORD = "docx"
    WORD_LEGACY = "doc"
    EXCEL = "xlsx"
    EXCEL_LEGACY = "xls"
    POWERPOINT = "pptx"
    POWERPOINT_LEGACY = "ppt"


class TranslationService:
    """Base class for translation services"""
    
    def __init__(self, source_lang='en', target_lang='zh-CN', api_key=None):
        self.source_lang = source_lang
        self.target_lang = target_lang
        self.api_key = api_key
        self.cache = {}
    
    @lru_cache(maxsize=10000)
    def translate_cached(self, text: str) -> str:
        """Translate with caching to reduce API calls"""
        if not text or not text.strip():
            return text
        return self.translate(text)
    
    def translate(self, text: str) -> str:
        """Override in subclasses"""
        raise NotImplementedError


class GoogleTranslateService(TranslationService):
    """Google Translate API wrapper"""
    
    def __init__(self, source_lang='en', target_lang='zh-CN', **kwargs):
        super().__init__(source_lang, target_lang)
        self.translator = GoogleTranslator(source=source_lang, target=target_lang)
    
    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
    def translate(self, text: str) -> str:
        if not text or not text.strip():
            return text
        try:
            return self.translator.translate(text)
        except Exception as e:
            print(f"Translation error: {e}")
            return text  # Return original on error


class DeepLTranslateService(TranslationService):
    """DeepL API wrapper"""
    
    def __init__(self, source_lang='en', target_lang='zh-CN', api_key=None, **kwargs):
        super().__init__(source_lang, target_lang, api_key)
        if not api_key:
            raise ValueError("DeepL requires an API key")
        
        if deepl:
            self.translator = deepl.Translator(api_key)
        else:
            # Fallback to deep-translator's DeepL
            self.translator = DeeplTranslator(
                api_key=api_key,
                source=source_lang,
                target=target_lang.split('-')[0].upper()
            )
    
    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
    def translate(self, text: str) -> str:
        if not text or not text.strip():
            return text
        try:
            if deepl and isinstance(self.translator, deepl.Translator):
                result = self.translator.translate_text(
                    text,
                    target_lang=self.target_lang.split('-')[0].upper()
                )
                return result.text
            else:
                return self.translator.translate(text)
        except Exception as e:
            print(f"Translation error: {e}")
            return text


class PDFTranslator:
    """PDF translation with format preservation using PyMuPDF"""
    
    def __init__(self, translation_service: TranslationService):
        if not pymupdf:
            raise ImportError("PyMuPDF is required for PDF translation")
        self.translator = translation_service
    
    def translate(self, input_path: str, output_path: str, preserve_original: bool = True) -> Dict:
        """
        Translate PDF while preserving formatting
        
        Args:
            input_path: Source PDF file path
            output_path: Destination PDF file path
            preserve_original: If True, create bilingual PDF with OCG layers
        
        Returns:
            Dictionary with translation statistics
        """
        doc = pymupdf.open(input_path)
        stats = {
            'total_pages': len(doc),
            'blocks_translated': 0,
            'blocks_failed': 0,
            'warnings': []
        }
        
        # Create optional content layer for translation if preserving original
        ocg_xref = None
        if preserve_original:
            ocg_xref = doc.add_ocg(f"{self.translator.target_lang} Translation", on=True)
        
        WHITE = pymupdf.pdfcolor["white"]
        
        for page_num, page in enumerate(doc):
            print(f"Processing PDF page {page_num + 1}/{len(doc)}...")
            
            # Extract text blocks with dehyphenation
            blocks = page.get_text("blocks", flags=pymupdf.TEXT_DEHYPHENATE)
            
            for block in blocks:
                if block[6] != 0:  # Skip non-text blocks
                    continue
                
                bbox = pymupdf.Rect(block[:4])
                original_text = block[4].strip()
                
                if not original_text:
                    continue
                
                try:
                    # Translate text
                    translated_text = self.translator.translate_cached(original_text)
                    
                    if preserve_original:
                        # Cover original with white rectangle (on translation layer)
                        page.draw_rect(bbox, color=None, fill=WHITE, oc=ocg_xref)
                    else:
                        # Redact original text entirely
                        page.add_redact_annot(bbox)
                    
                    # Insert translated text with automatic font selection
                    result = page.insert_htmlbox(
                        bbox,
                        translated_text,
                        css="""
                        * {
                            font-family: "Noto Sans CJK SC", "PingFang SC", "Microsoft YaHei", sans-serif;
                            font-size: 11pt;
                        }
                        """,
                        oc=ocg_xref if preserve_original else None,
                        scale_low=0.5  # Minimum 50% scaling
                    )
                    
                    if result < 0:
                        stats['warnings'].append(f"Page {page_num + 1}: Text overflow")
                    
                    stats['blocks_translated'] += 1
                    
                except Exception as e:
                    stats['blocks_failed'] += 1
                    stats['warnings'].append(f"Page {page_num + 1}: {str(e)}")
            
            if not preserve_original:
                page.apply_redactions()
        
        # Subset fonts to reduce file size
        doc.subset_fonts()
        
        # Save with compression
        doc.ez_save(output_path)
        doc.close()
        
        return stats


class ExcelTranslator:
    """Excel translation with formula preservation using openpyxl"""
    
    def __init__(self, translation_service: TranslationService):
        if not load_workbook:
            raise ImportError("openpyxl is required for Excel translation")
        self.translator = translation_service
        self.formula_tokenizer_available = Tokenizer is not None and Token is not None
    
    def translate(self, input_path: str, output_path: str) -> Dict:
        """
        Translate Excel while preserving formulas and formatting
        """
        wb = load_workbook(input_path, data_only=False)
        stats = {
            'sheets_processed': 0,
            'cells_translated': 0,
            'formulas_preserved': 0,
            'charts_translated': 0,
            'warnings': []
        }
        
        # Process each sheet
        for sheet_name in wb.sheetnames:
            print(f"Processing Excel sheet: {sheet_name}")
            sheet = wb[sheet_name]
            stats['sheets_processed'] += 1
            
            # Translate sheet name (optional)
            # new_name = self.translator.translate_cached(sheet_name)
            # sheet.title = new_name[:31]  # Excel limit is 31 chars
            
            # Process cells
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.data_type == 'f':  # Formula
                        if self.formula_tokenizer_available:
                            self._translate_formula(cell, stats)
                        # Whether or not we tokenized, formulas are preserved
                        stats['formulas_preserved'] += 1
                    elif isinstance(cell.value, str):
                        try:
                            cell.value = self.translator.translate_cached(cell.value)
                            stats['cells_translated'] += 1
                        except Exception as e:
                            stats['warnings'].append(f"Cell {cell.coordinate}: {str(e)}")
            
            # Translate charts
            if hasattr(sheet, '_charts'):
                for chart in sheet._charts:
                    self._translate_chart(chart, stats)
        
        # Save workbook
        wb.save(output_path)
        wb.close()
        
        return stats
    
    def _translate_formula(self, cell, stats):
        """Extract and translate text from formulas while preserving structure"""
        try:
            if Tokenizer is None or Token is None:
                # Tokenizer not available; leave formula as-is
                return
            formula = cell.value
            if not formula or not isinstance(formula, str):
                return
            tok = Tokenizer(formula)
            result_parts = []
            for token in tok.items:
                if token.type == Token.OPERAND and token.subtype == Token.TEXT:
                    text = token.value.strip('"')
                    if text:
                        translated = self.translator.translate_cached(text)
                        result_parts.append(f'"{translated}"')
                    else:
                        result_parts.append(token.value)
                else:
                    result_parts.append(token.value)
            cell.value = ''.join(result_parts)
        except Exception as e:
            stats['warnings'].append(f"Formula translation error: {str(e)}")
    
    def _translate_chart(self, chart, stats):
        """Translate chart titles and labels"""
        try:
            if chart.title and isinstance(chart.title, str):
                chart.title = self.translator.translate_cached(chart.title)
                stats['charts_translated'] += 1
            
            if hasattr(chart, 'x_axis') and chart.x_axis.title:
                chart.x_axis.title = self.translator.translate_cached(chart.x_axis.title)
            
            if hasattr(chart, 'y_axis') and chart.y_axis.title:
                chart.y_axis.title = self.translator.translate_cached(chart.y_axis.title)
        except Exception as e:
            stats['warnings'].append(f"Chart translation error: {str(e)}")


class WordTranslator:
    """Word translation using python-docx"""
    
    def __init__(self, translation_service: TranslationService):
        if not Document:
            raise ImportError("python-docx is required for Word translation")
        self.translator = translation_service
    
    def translate(self, input_path: str, output_path: str) -> Dict:
        """
        Translate Word document (with limitations for complex features)
        """
        doc = Document(input_path)
        stats = {
            'paragraphs_translated': 0,
            'tables_translated': 0,
            'headers_footers_translated': 0,
            'warnings': []
        }
        
        # Note: python-docx limitations
        if self._has_tracked_changes(input_path):
            stats['warnings'].append("Document has tracked changes - these will be lost")
        
        # Translate paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                self._translate_paragraph(para, stats)
        
        # Translate tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if para.text.strip():
                            self._translate_paragraph(para, stats)
                            stats['tables_translated'] += 1
        
        # Translate headers and footers
        for section in doc.sections:
            # Header
            for para in section.header.paragraphs:
                if para.text.strip():
                    self._translate_paragraph(para, stats)
                    stats['headers_footers_translated'] += 1
            
            # Footer
            for para in section.footer.paragraphs:
                if para.text.strip():
                    self._translate_paragraph(para, stats)
                    stats['headers_footers_translated'] += 1
        
        # Save document
        doc.save(output_path)
        
        return stats
    
    def _translate_paragraph(self, paragraph, stats):
        """Translate paragraph while preserving run-level formatting"""
        try:
            for run in paragraph.runs:
                if run.text.strip():
                    original_text = run.text
                    # Preserve leading/trailing whitespace
                    translated = self.translator.translate_cached(run.text.strip())
                    
                    if original_text.startswith(' '):
                        translated = ' ' + translated
                    if original_text.endswith(' '):
                        translated = translated + ' '
                    
                    run.text = translated
            stats['paragraphs_translated'] += 1
        except Exception as e:
            stats['warnings'].append(f"Paragraph translation error: {str(e)}")
    
    def _has_tracked_changes(self, filepath):
        """Check if document has tracked changes (basic check)"""
        # This is a simplified check - full detection would require XML parsing
        return False  # python-docx can't properly detect this


class PowerPointTranslator:
    """PowerPoint translation using python-pptx"""
    
    def __init__(self, translation_service: TranslationService):
        if not Presentation:
            raise ImportError("python-pptx is required for PowerPoint translation")
        self.translator = translation_service
    
    def translate(self, input_path: str, output_path: str) -> Dict:
        """
        Translate PowerPoint while preserving formatting and layouts
        """
        prs = Presentation(input_path)
        stats = {
            'slides_processed': len(prs.slides),
            'text_runs_translated': 0,
            'notes_translated': 0,
            'warnings': []
        }
        
        for slide_num, slide in enumerate(prs.slides):
            print(f"Processing PowerPoint slide {slide_num + 1}/{len(prs.slides)}...")
            
            # Process all shapes
            for shape in slide.shapes:
                self._translate_shape(shape, slide_num + 1, stats)
            
            # Process notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                self._translate_text_frame(notes_frame, slide_num + 1, stats, "notes")
                stats['notes_translated'] += 1
        
        prs.save(output_path)
        
        return stats
    
    def _translate_shape(self, shape, slide_num: int, stats: Dict):
        """Recursively translate shape text (handles groups)"""
        # Handle regular shapes with text
        if hasattr(shape, "text_frame") and shape.text_frame:
            self._translate_text_frame(shape.text_frame, slide_num, stats)
        
        # Handle tables
        if hasattr(shape, "table"):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        self._translate_text_frame(cell.text_frame, slide_num, stats, "table")
        
        # Handle grouped shapes (recursive)
        if hasattr(shape, 'shapes'):
            for sub_shape in shape.shapes:
                self._translate_shape(sub_shape, slide_num, stats)
    
    def _translate_text_frame(self, text_frame, slide_num: int, stats: Dict, context: str = "text"):
        """Translate all text in a text frame while preserving formatting"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text.strip()
                
                if not original_text:
                    continue
                
                try:
                    translated_text = self.translator.translate_cached(original_text)
                    
                    # Preserve leading/trailing whitespace
                    if run.text.startswith(' '):
                        translated_text = ' ' + translated_text
                    if run.text.endswith(' '):
                        translated_text = translated_text + ' '
                    
                    run.text = translated_text
                    stats['text_runs_translated'] += 1
                    
                except Exception as e:
                    stats['warnings'].append(f"Slide {slide_num} ({context}): {str(e)}")


class DocumentTranslator:
    """
    Main document translator that handles all supported formats
    """
    
    def __init__(self, 
                 source_lang: str = 'en',
                 target_lang: str = 'zh-CN',
                 translation_service: str = 'google',
                 api_key: Optional[str] = None):
        """
        Initialize document translator
        
        Args:
            source_lang: Source language code (e.g., 'en', 'zh-CN')
            target_lang: Target language code
            translation_service: 'google', 'deepl', or 'microsoft'
            api_key: API key for translation service (if required)
        """
        # Initialize translation service
        if translation_service == 'google':
            self.translation_service = GoogleTranslateService(source_lang, target_lang)
        elif translation_service == 'deepl':
            self.translation_service = DeepLTranslateService(source_lang, target_lang, api_key)
        else:
            raise ValueError(f"Unsupported translation service: {translation_service}")
        
        # Initialize format-specific translators
        self.pdf_translator = PDFTranslator(self.translation_service) if pymupdf else None
        self.excel_translator = ExcelTranslator(self.translation_service) if load_workbook else None
        self.word_translator = WordTranslator(self.translation_service) if Document else None
        self.ppt_translator = PowerPointTranslator(self.translation_service) if Presentation else None
    
    def translate_document(self, 
                          input_path: str, 
                          output_path: Optional[str] = None,
                          preserve_original: bool = True) -> Dict:
        """
        Auto-detect document type and translate accordingly
        
        Args:
            input_path: Source document path
            output_path: Destination path (auto-generated if None)
            preserve_original: For PDFs, whether to keep original as layer
            
        Returns:
            Translation statistics dictionary
        """
        input_file = Path(input_path)
        
        if not input_file.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")
        
        # Generate output path if not provided
        if output_path is None:
            output_path = str(input_file.parent / 
                            f"{input_file.stem}_translated{input_file.suffix}")
        
        # Detect file type and translate
        suffix = input_file.suffix.lower()
        
        print(f"\n{'='*60}")
        print(f"Translating: {input_file.name}")
        print(f"Format: {suffix}")
        print(f"Output: {Path(output_path).name}")
        print(f"{'='*60}\n")
        
        try:
            if suffix == '.pdf':
                if not self.pdf_translator:
                    raise ImportError("PyMuPDF required for PDF translation")
                stats = self.pdf_translator.translate(input_path, output_path, preserve_original)
                
            elif suffix == '.docx':
                if not self.word_translator:
                    raise ImportError("python-docx required for Word translation")
                stats = self.word_translator.translate(input_path, output_path)
                
            elif suffix in ['.xlsx', '.xlsm']:
                if not self.excel_translator:
                    raise ImportError("openpyxl required for Excel translation")
                stats = self.excel_translator.translate(input_path, output_path)
                
            elif suffix in ['.pptx', '.ppt']:
                if suffix == '.ppt':
                    print("⚠️  Warning: .ppt format not fully supported. Convert to .pptx first.")
                if not self.ppt_translator:
                    raise ImportError("python-pptx required for PowerPoint translation")
                stats = self.ppt_translator.translate(input_path, output_path)
                
            elif suffix == '.doc':
                raise ValueError("Legacy .doc format not supported. Please save as .docx")
                
            elif suffix == '.xls':
                raise ValueError("Legacy .xls format not supported. Please save as .xlsx")
                
            else:
                raise ValueError(f"Unsupported file format: {suffix}")
            
            # Print summary
            self._print_summary(stats)
            
            print(f"\n✅ Translation complete!")
            print(f"Output saved to: {output_path}")
            
            return stats
            
        except Exception as e:
            print(f"\n❌ Translation failed: {str(e)}")
            # Clean up partial output
            if Path(output_path).exists():
                Path(output_path).unlink()
            raise
    
    def translate_batch(self, 
                       input_folder: str,
                       output_folder: Optional[str] = None,
                       file_patterns: List[str] = None) -> Dict:
        """
        Translate multiple documents in a folder
        
        Args:
            input_folder: Source folder path
            output_folder: Destination folder (creates 'translated' subfolder if None)
            file_patterns: List of patterns like ['*.pdf', '*.docx']
        """
        input_path = Path(input_folder)
        
        if not input_path.exists():
            raise FileNotFoundError(f"Input folder not found: {input_folder}")
        
        # Set up output folder
        if output_folder is None:
            output_path = input_path / "translated"
        else:
            output_path = Path(output_folder)
        
        output_path.mkdir(parents=True, exist_ok=True)
        
        # Default patterns
        if file_patterns is None:
            file_patterns = ['*.pdf', '*.docx', '*.xlsx', '*.pptx']
        
        # Collect files
        files_to_translate = []
        for pattern in file_patterns:
            files_to_translate.extend(input_path.glob(pattern))
        
        print(f"Found {len(files_to_translate)} files to translate")
        
        # Process each file
        results = {
            'succeeded': [],
            'failed': [],
            'statistics': {}
        }
        
        for file_path in files_to_translate:
            output_file = output_path / f"{file_path.stem}_translated{file_path.suffix}"
            
            try:
                stats = self.translate_document(str(file_path), str(output_file))
                results['succeeded'].append(str(file_path))
                results['statistics'][str(file_path)] = stats
            except Exception as e:
                print(f"Failed to translate {file_path}: {str(e)}")
                results['failed'].append({'file': str(file_path), 'error': str(e)})
        
        # Summary
        print(f"\n{'='*60}")
        print(f"Batch Translation Summary")
        print(f"{'='*60}")
        print(f"✅ Succeeded: {len(results['succeeded'])} files")
        print(f"❌ Failed: {len(results['failed'])} files")
        
        return results
    
    def _print_summary(self, stats: Dict):
        """Print translation statistics summary"""
        print("\n📊 Translation Statistics:")
        print("-" * 40)
        
        for key, value in stats.items():
            if key != 'warnings':
                print(f"  {key}: {value}")
        
        if stats.get('warnings'):
            print(f"\n⚠️  Warnings ({len(stats['warnings'])}):")
            for warning in stats['warnings'][:5]:  # Show first 5 warnings
                print(f"  - {warning}")
            if len(stats['warnings']) > 5:
                print(f"  ... and {len(stats['warnings']) - 5} more")


def main():
    """Command-line interface"""
    parser = argparse.ArgumentParser(
        description='Translate documents while preserving formatting'
    )
    
    parser.add_argument('input', help='Input file or folder path')
    parser.add_argument('-o', '--output', help='Output file or folder path')
    parser.add_argument('-s', '--source', default='en', help='Source language (default: en)')
    parser.add_argument('-t', '--target', default='zh-CN', help='Target language (default: zh-CN)')
    parser.add_argument('--service', default='google', 
                       choices=['google', 'deepl'],
                       help='Translation service (default: google)')
    parser.add_argument('--api-key', help='API key for translation service')
    parser.add_argument('--batch', action='store_true', help='Batch process folder')
    parser.add_argument('--no-preserve', action='store_true', 
                       help="Don't preserve original in PDFs")
    
    args = parser.parse_args()
    
    # Create translator
    translator = DocumentTranslator(
        source_lang=args.source,
        target_lang=args.target,
        translation_service=args.service,
        api_key=args.api_key
    )
    
    # Process
    if args.batch:
        results = translator.translate_batch(args.input, args.output)
        
        # Save results to JSON
        results_file = Path(args.output or args.input) / "translation_results.json"
        with open(results_file, 'w', encoding='utf-8') as f:
            json.dump(results, f, indent=2, ensure_ascii=False)
        print(f"\nResults saved to: {results_file}")
    else:
        translator.translate_document(
            args.input, 
            args.output,
            preserve_original=not args.no_preserve
        )


if __name__ == "__main__":
    # Check for required libraries
    required = []
    if not pymupdf:
        required.append("pymupdf")
    if not load_workbook:
        required.append("openpyxl")
    if not Document:
        required.append("python-docx")
    if not Presentation:
        required.append("python-pptx")
    
    if required:
        print("\n⚠️  Missing required libraries. Install with:")
        print(f"pip install {' '.join(required)} deep-translator tenacity")
        print("\nFor full functionality, install all dependencies:")
        print("pip install pymupdf openpyxl python-docx python-pptx deep-translator tenacity")
        print("\nOptional for better translation:")
        print("pip install deepl  # For DeepL API (requires API key)")
        sys.exit(1)
    
    main()